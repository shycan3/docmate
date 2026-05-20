from pptx import Presentation
from pptx.util import Inches
import os

prs = Presentation()
img_dir = os.path.join(os.getcwd(), "presentation-materials", "images")
out_path = os.path.join(os.getcwd(), "presentation-materials", "DocMate_Presentation_초안.pptx")

# Title slide
slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]
title.text = "DocMate — 청년 정책 통합 플랫폼"
subtitle.text = "경영정보시스템 프로젝트 경진대회 발표\n팀명 · 소속 · 연락처"

# Insert poster image on title slide right side
img_path = os.path.join(img_dir, "01-first-screen-analysis-workspace.png")
if os.path.exists(img_path):
    try:
        slide.shapes.add_picture(img_path, Inches(5.5), Inches(1.0), width=Inches(4.0))
    except Exception as e:
        print('Could not insert title image:', e)

# helper to add bullet slides and optional image
def add_bullet_slide(title_text, bullets, img_name=None, img_pos=(4.5,1.5), img_w=4.0):
    layout = prs.slide_layouts[1]
    s = prs.slides.add_slide(layout)
    s.shapes.title.text = title_text
    body = s.shapes.placeholders[1].text_frame
    # clear default
    body.clear()
    for i, b in enumerate(bullets):
        if i == 0:
            p = body.paragraphs[0]
            p.text = b
        else:
            p = body.add_paragraph()
            p.text = b
        p.level = 0
    if img_name:
        pth = os.path.join(img_dir, img_name)
        if os.path.exists(pth):
            try:
                s.shapes.add_picture(pth, Inches(img_pos[0]), Inches(img_pos[1]), width=Inches(img_w))
            except Exception as e:
                print('Could not insert image', img_name, e)

# Slides based on plan
add_bullet_slide("문제 정의", [
    "청년이 정책을 '찾고·이해하고·신청'하기 어려움",
    "분산된 정보, 복잡한 자격요건, 낮은 체감성(정책 피로)"
])

add_bullet_slide("시장 데이터(요약)", [
    "'온통청년' 인지도 55.1% / '내용 잘 안다' 7.8% (정부, 2024)",
    "청년 885명 조사: 이용경험 83.6% / 운영 만족 26.4% (한대신문, 2025)"
])

add_bullet_slide("증거 및 구조적 원인", [
    "부처별 파편화로 개인 자가진단 어려움 (KDI, 2025)"
], img_name="04-source-evidence-expanded.png", img_pos=(4.5,1.8))

add_bullet_slide("페르소나 & 여정", [
    "A: 대학생(취업준비) — 정보탐색 → 자가진단 → 신청"
], img_name="02-structured-profile-and-sample-choice.png", img_pos=(4.5,1.5))

add_bullet_slide("솔루션(요약)", [
    "통합검색 + 개인정보 기반 자가진단",
    "간단 모드: 핵심 3~4개 입력"
], img_name="01-first-screen-analysis-workspace.png", img_pos=(4.5,1.5))

add_bullet_slide("데모 흐름", [
    "입력 → 결과(우선순위) → 신청 체크리스트"
], img_name="03-result-summary-decision.png", img_pos=(4.5,1.5))

add_bullet_slide("아키텍처 & 개인정보", [
    "데이터 수집(정부 API/크롤링) → 정규화 → 매칭",
    "개인정보: 로컬 처리·익명화, 동의 필요"
], img_name="06-history-comparison-demo-mode.png", img_pos=(4.5,1.5))

add_bullet_slide("Go-To-Market", [
    "대학·취업센터 파일럿, 지자체 제휴, SNS 캠페인"
])

add_bullet_slide("KPI & 측정", [
    "DAU/MAU, 신청전환율(신청완료/권유받은수)",
    "자가진단 성공률, 만족도 개선목표"
])

add_bullet_slide("로드맵 & 리스크", [
    "MVP(3개월) → 파일럿(3개월) → 확장(6개월)",
    "리스크: 데이터제휴/법률/채택부족 — 완화책 제시"
])

# Final slide
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = "요약 및 Q&A"
tf = slide.shapes.placeholders[1].text_frame
tf.text = "요약: 통합·간편 자가진단으로 청년의 정책 접근성 개선\n문의: 팀명 · 연락처"

try:
    prs.save(out_path)
    print('Saved:', out_path)
except Exception as e:
    print('Failed to save presentation:', e)
